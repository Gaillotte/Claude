#!/usr/bin/env python3
"""
Génère le deck de slides didactiques AI Transit Pipeline.
Export : .pptx  (puis converti en PDF via LibreOffice)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm
from pptx.oxml.ns import qn
from lxml import etree
import copy, textwrap

# ── Palette ──────────────────────────────────────────────────────────────────
C_BG_DARK   = RGBColor(0x0D, 0x1B, 0x2A)   # fond principal bleu nuit
C_BG_CARD   = RGBColor(0x1B, 0x2A, 0x3B)   # fond carte
C_ACCENT    = RGBColor(0x00, 0xA8, 0xFF)   # bleu électrique
C_ACCENT2   = RGBColor(0x00, 0xD4, 0xAA)   # vert cyan
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xCC, 0xDD, 0xEE)
C_GRAY      = RGBColor(0x88, 0x99, 0xAA)
C_GREEN     = RGBColor(0x2E, 0xCC, 0x71)
C_RED       = RGBColor(0xE7, 0x4C, 0x3C)
C_ORANGE    = RGBColor(0xF3, 0x9C, 0x12)
C_YELLOW    = RGBColor(0xF1, 0xC4, 0x0F)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank_layout = prs.slide_layouts[6]   # totalement vide


# ═══════════════════════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════════════════════
def rgb_hex(c: RGBColor) -> str:
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"

def add_slide():
    return prs.slides.add_slide(blank_layout)

def fill_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, x, y, w, h, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE=1
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def rounded_box(slide, x, y, w, h, color: RGBColor, radius=Pt(8)):
    shape = slide.shapes.add_shape(5, x, y, w, h)   # ROUNDED_RECTANGLE=5
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    adj = shape.adjustments
    adj[0] = 0.05
    return shape

def txt(slide, text, x, y, w, h,
        size=24, bold=False, italic=False, color=C_WHITE,
        align=PP_ALIGN.LEFT, wrap=True, font="Calibri"):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name   = font
    return txb

def txt_lines(slide, lines, x, y, w, h,
              size=18, bold=False, color=C_WHITE,
              align=PP_ALIGN.LEFT, spacing=1.2):
    """Ajoute plusieurs lignes dans une même zone de texte."""
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(size)
        run.font.bold  = bold
        run.font.color.rgb = color
        run.font.name  = "Calibri"
        p.space_after  = Pt(4)
    return txb

def divider(slide, y, color=C_ACCENT):
    box(slide, Inches(0.5), y, W - Inches(1), Pt(2), color)

def pill(slide, text, x, y, color: RGBColor, text_color=C_WHITE, size=13):
    w = Inches(2.2)
    h = Inches(0.38)
    s = slide.shapes.add_shape(5, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.adjustments[0] = 0.5
    tf = s.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = True
    run.font.color.rgb = text_color
    run.font.name  = "Calibri"
    return s

def mono_box(slide, text, x, y, w, h, size=10):
    """Boîte style code monospace."""
    bg = slide.shapes.add_shape(1, x, y, w, h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x06, 0x0F, 0x18)
    bg.line.color.rgb = C_ACCENT
    bg.line.width = Pt(1)

    txb = slide.shapes.add_textbox(
        x + Inches(0.12), y + Inches(0.1),
        w - Inches(0.24), h - Inches(0.2))
    tf = txb.text_frame
    tf.word_wrap = False
    first = True
    for line in text.split("\n"):
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size  = Pt(size)
        run.font.name  = "Courier New"
        run.font.color.rgb = RGBColor(0x7F, 0xFF, 0xD4)
    return txb

def badge(slide, text, x, y, color):
    s = slide.shapes.add_shape(5, x, y, Inches(0.9), Inches(0.3))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.adjustments[0] = 0.5
    tf = s.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = text
    r.font.size  = Pt(10)
    r.font.bold  = True
    r.font.color.rgb = C_WHITE
    r.font.name  = "Calibri"

def section_header(slide, title, subtitle=""):
    fill_bg(slide, C_BG_DARK)
    # Bande accent gauche
    box(slide, 0, 0, Inches(0.18), H, C_ACCENT)
    # Fond carte centré
    rounded_box(slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(3.2), C_BG_CARD)
    txt(slide, title,    Inches(2), Inches(2.6), Inches(9.3), Inches(1.2),
        size=40, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    if subtitle:
        txt(slide, subtitle, Inches(2), Inches(3.7), Inches(9.3), Inches(0.7),
            size=20, italic=True, color=C_LIGHT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COUVERTURE
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG_DARK)

# Gradient simulé : deux rectangles
box(sl, 0, 0, W, Inches(1.2), RGBColor(0x00, 0xA8, 0xFF))
box(sl, 0, Inches(6.3), W, Inches(1.2), RGBColor(0x00, 0xD4, 0xAA))

txt(sl, "AI TRANSIT PIPELINE",
    Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.4),
    size=52, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

txt(sl, "Sécuriser l'intégration du code IA en entreprise",
    Inches(1), Inches(2.9), Inches(11.3), Inches(0.8),
    size=24, italic=True, color=C_ACCENT2, align=PP_ALIGN.CENTER)

divider(sl, Inches(3.85), C_GRAY)

bullets_cover = [
    "🔍  Récupération sécurisée depuis GitHub",
    "🛡️   Scan multicouche adaptatif par type de fichier",
    "📦  Archive ZIP approuvée + rapport Excel de traçabilité",
]
txt_lines(sl, bullets_cover,
          Inches(3), Inches(4.1), Inches(7.3), Inches(2),
          size=18, color=C_LIGHT, align=PP_ALIGN.LEFT)

txt(sl, "Version 1.0  —  Juin 2026",
    Inches(0), Inches(7.0), W, Inches(0.4),
    size=11, color=C_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — SOMMAIRE
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG_DARK)
box(sl, 0, 0, W, Inches(1.1), RGBColor(0x00, 0x70, 0xA8))
txt(sl, "Sommaire", Inches(0.5), Inches(0.15), W, Inches(0.85),
    size=36, bold=True, color=C_WHITE)
divider(sl, Inches(1.15))

items = [
    ("01", "Contexte & objectifs",            C_ACCENT),
    ("02", "Architecture du pipeline",         C_ACCENT2),
    ("03", "Phase 1 — Récupération sécurisée", C_YELLOW),
    ("04", "Phase 2 — Scan multicouche",        C_ORANGE),
    ("05", "Détail des checks par type",        RGBColor(0xE0, 0x70, 0xFF)),
    ("06", "Outils utilisés",                   C_GREEN),
    ("07", "Résultats & livrables",             C_ACCENT),
    ("08", "Règles de sécurité absolues",       C_RED),
    ("09", "Exploitation & assurance qualité",   C_ACCENT2),
]

cols = [(Inches(0.5), Inches(6.7))]
for i, (num, title, color) in enumerate(items):
    # 9 items: 5 in the left column, 4 in the right, rows tightened to fit.
    col_x = Inches(0.5) if i < 5 else Inches(6.7)
    row   = i if i < 5 else i - 5
    y = Inches(1.35 + row * 1.12)

    rounded_box(sl, col_x, y, Inches(5.9), Inches(0.95), C_BG_CARD)
    # Numéro
    s = sl.shapes.add_shape(5, col_x + Inches(0.12), y + Inches(0.15),
                            Inches(0.65), Inches(0.65))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.adjustments[0] = 0.5
    tf = s.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run(); r.text = num
    r.font.size = Pt(14); r.font.bold = True
    r.font.color.rgb = C_WHITE; r.font.name = "Calibri"
    # Titre
    txt(sl, title, col_x + Inches(0.9), y + Inches(0.2),
        Inches(4.8), Inches(0.6), size=16, bold=True, color=C_WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — CONTEXTE
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG_DARK)
box(sl, 0, 0, W, Inches(1.1), C_ACCENT)
txt(sl, "01  Contexte & Objectifs", Inches(0.5), Inches(0.15), W, Inches(0.85),
    size=32, bold=True, color=C_WHITE)
divider(sl, Inches(1.15))

txt(sl, "Pourquoi ce pipeline ?",
    Inches(0.5), Inches(1.3), Inches(8), Inches(0.6),
    size=22, bold=True, color=C_ACCENT)

problems = [
    "🤖  Le code IA généré peut contenir des backdoors, secrets ou patterns dangereux",
    "📥  Les développeurs importent du code sans revue de sécurité systématique",
    "🏢  Le réseau d'entreprise doit rester isolé (air-gap partiel)",
    "📋  Aucune traçabilité sans processus structuré",
]
txt_lines(sl, problems, Inches(0.5), Inches(1.95), Inches(7.8), Inches(2.6),
          size=17, color=C_LIGHT)

# Schéma simple : Internet → Pipeline → Réseau interne
for label, x, color in [
    ("🌐 Internet\n(GitHub)", Inches(9.0), C_ACCENT),
    ("🛡️ Pipeline\nAI Transit", Inches(10.2), C_ORANGE),
    ("🏢 Réseau\ninterne", Inches(11.4), C_GREEN),
]:
    s = sl.shapes.add_shape(5, x, Inches(2.2), Inches(1.05), Inches(0.95))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.adjustments[0] = 0.08
    tf = s.text_frame; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for li, line in enumerate(label.split("\n")):
        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = line
        r.font.size = Pt(11); r.font.bold = True
        r.font.color.rgb = C_WHITE; r.font.name = "Calibri"

for ax in [Inches(10.08), Inches(11.28)]:
    txt(sl, "→", ax, Inches(2.5), Inches(0.3), Inches(0.4),
        size=20, bold=True, color=C_GRAY, align=PP_ALIGN.CENTER)

txt(sl, "Flux\nunidirectionnel", Inches(9.0), Inches(3.3), Inches(3.5), Inches(0.6),
    size=11, italic=True, color=C_GRAY, align=PP_ALIGN.CENTER)

# Objectifs
txt(sl, "Objectifs clés",
    Inches(0.5), Inches(4.6), Inches(12), Inches(0.5),
    size=20, bold=True, color=C_ACCENT2)

obj = [
    ("Récupérer", "Clone sécurisé depth 1", C_ACCENT),
    ("Scanner",   "Multicouche adaptatif",  C_ORANGE),
    ("Décider",   "PASS → ZIP  /  FAIL → Quarantaine", C_RED),
    ("Tracer",    "Rapport Excel + JSON + HTML", C_GREEN),
]
for i, (title, sub, color) in enumerate(obj):
    x = Inches(0.4 + i * 3.2)
    rounded_box(sl, x, Inches(5.1), Inches(3.0), Inches(1.85), C_BG_CARD)
    s2 = sl.shapes.add_shape(5, x + Inches(0.1), Inches(5.2),
                              Inches(0.5), Inches(0.4))
    s2.fill.solid(); s2.fill.fore_color.rgb = color
    s2.line.fill.background(); s2.adjustments[0] = 0.5
    txt(sl, title, x + Inches(0.7), Inches(5.2), Inches(2.1), Inches(0.45),
        size=16, bold=True, color=C_WHITE)
    txt(sl, sub, x + Inches(0.15), Inches(5.7), Inches(2.75), Inches(0.7),
        size=13, color=C_LIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG_DARK)
box(sl, 0, 0, W, Inches(1.1), C_ACCENT2)
txt(sl, "02  Architecture du pipeline", Inches(0.5), Inches(0.15), W, Inches(0.85),
    size=32, bold=True, color=C_BG_DARK)
divider(sl, Inches(1.15), C_ACCENT2)

# Diagramme flux vertical simplifié
nodes = [
    ("📂 Source\nGitHub / Local",  Inches(0.5),  Inches(1.3),  C_ACCENT),
    ("⬇  fetch_repo.sh\nPhase 1",  Inches(0.5),  Inches(2.35), RGBColor(0x25, 0x55, 0x85)),
    ("🔍 scan_pipeline.sh\nPhase 2",Inches(0.5),  Inches(3.4),  RGBColor(0x85, 0x45, 0x10)),
    ("✔ Good/ (ZIP + Excel)",       Inches(0.5),  Inches(4.95), C_GREEN),
    ("✘ quarantine/ (chmod 700)",  Inches(0.5),  Inches(5.9),  C_RED),
]
for label, x, y, color in nodes:
    rounded_box(sl, x, y, Inches(3.4), Inches(0.9), color)
    for li, line in enumerate(label.split("\n")):
        p_y = y + Inches(0.07 + li * 0.38)
        txt(sl, line, x + Inches(0.1), p_y, Inches(3.2), Inches(0.42),
            size=14, bold=(li==0), color=C_WHITE)

# Flèches
for y_a in [Inches(2.25), Inches(3.3)]:
    txt(sl, "▼", Inches(1.8), y_a, Inches(0.6), Inches(0.4),
        size=18, color=C_GRAY, align=PP_ALIGN.CENTER)

txt(sl, "PASS ▼", Inches(0.7), Inches(4.85), Inches(1.5), Inches(0.35),
    size=13, color=C_GREEN, bold=True)
txt(sl, "FAIL ▼", Inches(2.2), Inches(4.85), Inches(1.5), Inches(0.35),
    size=13, color=C_RED, bold=True)

# Détail Phase 1
txt(sl, "Phase 1 — fetch_repo.sh",
    Inches(4.2), Inches(1.3), Inches(4.5), Inches(0.45),
    size=16, bold=True, color=C_ACCENT)
p1_items = [
    "✦  Whitelist hôtes (github.com uniquement)",
    "✦  Vérif. taille via API GitHub (< 500 MB)",
    "✦  git clone --depth 1 --no-tags",
    "✦  Suppression .git/ (pas de metadata)",
    "✦  Manifest SHA-256 de tous les fichiers",
    "✦  Triage rapide patterns suspects",
]
txt_lines(sl, p1_items, Inches(4.2), Inches(1.8), Inches(4.3), Inches(2.0),
          size=13, color=C_LIGHT)

# Détail Phase 2
txt(sl, "Phase 2 — scan_pipeline.sh",
    Inches(8.9), Inches(1.3), Inches(4.3), Inches(0.45),
    size=16, bold=True, color=C_ORANGE)
p2_items = [
    "🌐  Couche GLOBALE :",
    "     Gitleaks · detect-secrets · ClamAV · YARA",
    "",
    "📄  Couche PAR TYPE :",
    "     .py → Bandit + eval/exec",
    "     .js/.ts → Semgrep + child_process",
    "     .sh → ShellCheck + curl|bash",
    "     .yml → actions non-pinnées",
    "     .tf → checkov + secrets hardcodés",
    "     Dockerfile → hadolint + :latest",
    "     .so/.exe → FAIL auto + strings",
]
txt_lines(sl, p2_items, Inches(8.9), Inches(1.8), Inches(4.3), Inches(3.5),
          size=12, color=C_LIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — PHASE 1 FETCH
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG_DARK)
box(sl, 0, 0, W, Inches(1.1), C_YELLOW)
txt(sl, "03  Phase 1 — Récupération sécurisée", Inches(0.5), Inches(0.15), W, Inches(0.85),
    size=30, bold=True, color=C_BG_DARK)
divider(sl, Inches(1.15), C_YELLOW)

steps = [
    ("1", "Whitelist hôtes",
     "Seul github.com est autorisé. Toute autre URL est rejetée immédiatement.",
     C_ACCENT),
    ("2", "Vérification taille",
     "Appel API GitHub : taille du repo vérifiée avant tout clone. Limite : 500 MB.",
     C_ACCENT2),
    ("3", "Clone minimal",
     "git clone --depth 1 --no-tags --single-branch\nHistorique zéro, surface d'attaque minimale.",
     C_YELLOW),
    ("4", "Suppression .git/",
     "Les métadonnées Git (hooks, remotes, submodules) sont entièrement supprimées.",
     C_ORANGE),
    ("5", "Manifest SHA-256",
     "Hash SHA-256 de chaque fichier → fichier .manifest_sha256.txt pour audit d'intégrité.",
     RGBColor(0xE0, 0x70, 0xFF)),
    ("6", "Triage rapide",
     "Grep sur eval(, exec(, curl|bash, rm -rf → alerte immédiate avant scan approfondi.",
     C_GREEN),
]

for i, (num, title, desc, color) in enumerate(steps):
    col = i % 3
    row = i // 3
    x = Inches(0.4 + col * 4.3)
    y = Inches(1.35 + row * 2.85)
    rounded_box(sl, x, y, Inches(4.1), Inches(2.6), C_BG_CARD)
    # Numéro badge
    s = sl.shapes.add_shape(5, x + Inches(0.12), y + Inches(0.12),
                            Inches(0.5), Inches(0.5))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.adjustments[0] = 0.5
    tf2 = s.text_frame; tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    r2 = tf2.paragraphs[0].add_run(); r2.text = num
    r2.font.size = Pt(14); r2.font.bold = True
    r2.font.color.rgb = C_BG_DARK; r2.font.name = "Calibri"
    txt(sl, title, x + Inches(0.75), y + Inches(0.15),
        Inches(3.2), Inches(0.5), size=15, bold=True, color=color)
    txt(sl, desc, x + Inches(0.15), y + Inches(0.75),
        Inches(3.8), Inches(1.6), size=13, color=C_LIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PHASE 2 COUCHE GLOBALE
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG_DARK)
box(sl, 0, 0, W, Inches(1.1), C_ORANGE)
txt(sl, "04  Phase 2 — Couche globale", Inches(0.5), Inches(0.15), W, Inches(0.85),
    size=32, bold=True, color=C_WHITE)
divider(sl, Inches(1.15), C_ORANGE)

txt(sl, "S'applique à TOUT le répertoire, quel que soit le type de fichier",
    Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.5),
    size=17, italic=True, color=C_LIGHT, align=PP_ALIGN.CENTER)

tools_global = [
    ("Gitleaks", "🔑",
     ["Tokens GitHub / GitLab / AWS / GCP / Azure",
      "Clés privées RSA, PEM, SSH",
      "Clés API Stripe, Twilio, SendGrid, JWT",
      "> 150 types de credentials couverts",
      "Analyse entropique + règles regex"],
     C_ACCENT),
    ("detect-secrets", "📊",
     ["Entropie de Shannon sur les chaînes",
      "Détecte les secrets inconnus des règles",
      "Base64 et hex à haute densité aléatoire",
      "Complémentaire à Gitleaks",
      "Outil Yelp (battle-tested en production)"],
     C_ACCENT2),
    ("ClamAV", "🦠",
     ["Millions de signatures malware connues",
      "Trojans, ransomwares, backdoors, exploits",
      "Scan récursif de tous les fichiers",
      "Base mise à jour via freshclam",
      "Détection de binaires embarqués"],
     C_ORANGE),
    ("YARA", "🎯",
     ["Standard industriel IOC",
      "Règles .yar personnalisées dans yara-rules/",
      "Patterns propres aux backdoors IA",
      "Combinaison : regex + logique booléenne",
      "Extensible à volonté"],
     RGBColor(0xE0, 0x70, 0xFF)),
]

for i, (name, icon, items, color) in enumerate(tools_global):
    x = Inches(0.35 + i * 3.25)
    rounded_box(sl, x, Inches(1.85), Inches(3.05), Inches(5.35), C_BG_CARD)
    # Header outil
    s = sl.shapes.add_shape(1, x, Inches(1.85), Inches(3.05), Inches(0.65))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    txt(sl, f"{icon}  {name}", x + Inches(0.08), Inches(1.9),
        Inches(2.9), Inches(0.55), size=15, bold=True,
        color=C_BG_DARK if name == "detect-secrets" else C_WHITE,
        align=PP_ALIGN.CENTER)
    # Bullet items
    for j, item in enumerate(items):
        txt(sl, f"• {item}", x + Inches(0.1), Inches(2.65 + j * 0.75),
            Inches(2.85), Inches(0.65), size=12, color=C_LIGHT)

    badge(sl, "FAIL si positif", x + Inches(0.9), Inches(7.05), C_RED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — COUCHE PAR TYPE (tableau)
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG_DARK)
box(sl, 0, 0, W, Inches(1.1), RGBColor(0xE0, 0x70, 0xFF))
txt(sl, "04  Phase 2 — Dispatch par type de fichier",
    Inches(0.5), Inches(0.15), W, Inches(0.85),
    size=30, bold=True, color=C_WHITE)
divider(sl, Inches(1.15), RGBColor(0xE0, 0x70, 0xFF))

headers = ["Extension(s)", "Outil SAST", "Checks manuels", "Verdict défaut"]
col_x   = [Inches(0.3), Inches(2.15), Inches(4.75), Inches(10.85)]
col_w   = [Inches(1.75), Inches(2.45), Inches(6.0),  Inches(2.0)]

# En-têtes
for i, (hdr, x, w) in enumerate(zip(headers, col_x, col_w)):
    s = sl.shapes.add_shape(1, x, Inches(1.25), w, Inches(0.48))
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor(0x60, 0x20, 0xA0)
    s.line.fill.background()
    txt(sl, hdr, x + Inches(0.05), Inches(1.28), w - Inches(0.1), Inches(0.4),
        size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

rows_type = [
    (".py",             "Bandit",      "eval() exec() import os/subprocess",           "FAIL"),
    (".js .ts .jsx",    "Semgrep",     "eval() child_process require(…)",               "FAIL"),
    (".c .cpp .h",      "cppcheck",    "gets() strcpy() system() popen()",              "FAIL"),
    (".sh .bash .ps1",  "ShellCheck",  "curl|bash  wget|sh  eval $var",                 "FAIL"),
    (".yml .yaml",      "—",           "Actions non-pinnées  secrets inline",           "FAIL"),
    (".tf .tfvars",     "checkov",     "Secrets en dur  IAM trop larges",               "FAIL"),
    ("Dockerfile",      "hadolint",    ":latest  ADD http://  RUN curl|bash",           "FAIL"),
    (".so .exe .elf",   "strings",     "FAIL automatique + recherche IOC",              "FAIL ⚠"),
    (".zip .tar.gz",    "—",           "FAIL auto — extraction + re-scan requis",       "FAIL ⚠"),
    (".sql",            "—",           "xp_cmdshell  DROP TABLE  injection",            "FAIL"),
    (".json .md .txt",  "—",           "password: secret: token: api_key=",             "FAIL"),
]

for ri, (ext, tool, checks, verdict) in enumerate(rows_type):
    y = Inches(1.77 + ri * 0.5)
    bg = RGBColor(0x12, 0x22, 0x33) if ri % 2 == 0 else C_BG_CARD
    for ci, (val, x, w) in enumerate(zip([ext, tool, checks, verdict],
                                          col_x, col_w)):
        s = sl.shapes.add_shape(1, x, y, w, Inches(0.46))
        s.fill.solid(); s.fill.fore_color.rgb = bg; s.line.fill.background()
        color = C_RED if ci == 3 else (C_ACCENT if ci == 1 else C_LIGHT)
        txt(sl, val, x + Inches(0.05), y + Inches(0.05),
            w - Inches(0.1), Inches(0.36),
            size=11, color=color,
            bold=(ci == 0 or ci == 3))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES 8–10 — DÉTAIL CHECKS PAR TYPE (3 slides regroupés)
# ═══════════════════════════════════════════════════════════════════════════════
details = [
    # (type_name, color, desc, checks)
    ("Python .py", C_ACCENT,
     "Langage à exécution dynamique — risque élevé de backdoors via eval/exec",
     [
         ("Bandit SAST",       "Analyse AST Python : subprocess shell=True, pickle, MD5, SQL injection, assert sécurité…", "FAIL"),
         ("eval() / exec()",   "Exécution arbitraire de code à l'exécution — backdoor triviale dans du code IA généré",    "FAIL"),
         ("import os/subprocess", "Accès au shell système — légitime mais nécessite revue manuelle",                       "WARN"),
     ]),
    ("JavaScript/TS .js .ts", C_YELLOW,
     "Écosystème exposé aux supply-chain attacks et à l'exécution dynamique",
     [
         ("Semgrep p/javascript", "XSS, prototype pollution, innerHTML, open redirect, SSTI…",   "FAIL"),
         ("eval()",               "Exécution dynamique de code — vecteur classique d'injection", "FAIL"),
         ("child_process",        "Exécution de commandes shell depuis Node.js",                  "FAIL"),
     ]),
    ("C / C++ .c .cpp", C_ORANGE,
     "Risques liés à la gestion manuelle de la mémoire",
     [
         ("cppcheck",    "Buffer overflows, null ptr, use-after-free, mémoire non initialisée",         "FAIL"),
         ("gets()",      "Lecture sans limite de taille → buffer overflow garanti",                     "FAIL"),
         ("strcpy()",    "Copie sans vérification de taille → débordement",                            "FAIL"),
         ("system()/popen()", "Exécution de commandes shell → injection possible",                     "FAIL"),
     ]),
    ("Shell .sh .bash", C_ACCENT2,
     "Exécution directe — syntaxe propice à l'obfuscation de commandes",
     [
         ("ShellCheck",      "Variables non quotées, globbing, redirections dangereuses…",             "FAIL"),
         ("curl | bash",     "Exécution de code distant sans vérification d'intégrité",               "FAIL"),
         ("eval $variable",  "Injection triviale si la variable est contrôlée par un attaquant",       "FAIL"),
     ]),
    ("YAML / CI .yml", RGBColor(0xE0, 0x70, 0xFF),
     "Pipelines CI/CD — une mauvaise config ouvre des backdoors d'infrastructure",
     [
         ("Actions non-pinnées",  "uses: action@main — vulnérable au tag-hijacking si le tag est réécrit", "FAIL"),
         ("Secrets inline",       "password: valeur_en_clair (pas ${{ secrets.XXX }})",                   "FAIL"),
     ]),
    ("Terraform .tf", C_GREEN,
     "Code IaC — peut provisionner des ressources cloud entières",
     [
         ("checkov (CIS)",   "Buckets S3 sans chiffrement, IAM trop larges, BDD sans chiffrement at-rest", "FAIL"),
         ("Secrets en dur",  "password = \"valeur\", token = \"valeur\" dans les fichiers .tf/.tfvars",    "FAIL"),
     ]),
]

for group in range(0, len(details), 2):
    sl = add_slide()
    fill_bg(sl, C_BG_DARK)
    box(sl, 0, 0, W, Inches(1.1), RGBColor(0x60, 0x20, 0xA0))
    txt(sl, "05  Détail des vérifications par type", Inches(0.5), Inches(0.15),
        W, Inches(0.85), size=30, bold=True, color=C_WHITE)
    divider(sl, Inches(1.15), RGBColor(0xE0, 0x70, 0xFF))

    for col_idx in range(2):
        di = group + col_idx
        if di >= len(details):
            break
        name, color, desc, checks = details[di]
        cx = Inches(0.3) if col_idx == 0 else Inches(6.7)
        cw = Inches(6.2)

        rounded_box(sl, cx, Inches(1.25), cw, Inches(6.0), C_BG_CARD)
        # Header
        s = sl.shapes.add_shape(1, cx, Inches(1.25), cw, Inches(0.55))
        s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
        txt(sl, name, cx + Inches(0.1), Inches(1.29), cw - Inches(0.2), Inches(0.45),
            size=17, bold=True, color=C_BG_DARK)
        txt(sl, desc, cx + Inches(0.1), Inches(1.87), cw - Inches(0.2), Inches(0.55),
            size=12, italic=True, color=C_GRAY)

        for ci, (check, expl, verdict) in enumerate(checks):
            cy = Inches(2.55 + ci * 1.3)
            bg2 = RGBColor(0x0E, 0x1E, 0x2E) if ci % 2 == 0 else RGBColor(0x16, 0x2A, 0x3C)
            s2 = sl.shapes.add_shape(1, cx + Inches(0.1), cy,
                                     cw - Inches(0.2), Inches(1.18))
            s2.fill.solid(); s2.fill.fore_color.rgb = bg2; s2.line.fill.background()
            vcolor = C_RED if verdict == "FAIL" else C_ORANGE
            txt(sl, check, cx + Inches(0.2), cy + Inches(0.04),
                cw - Inches(1.1), Inches(0.38), size=13, bold=True, color=vcolor)
            badge_s = sl.shapes.add_shape(5,
                cx + cw - Inches(1.1), cy + Inches(0.07),
                Inches(0.9), Inches(0.3))
            badge_s.fill.solid(); badge_s.fill.fore_color.rgb = vcolor
            badge_s.line.fill.background(); badge_s.adjustments[0] = 0.5
            tf3 = badge_s.text_frame; tf3.paragraphs[0].alignment = PP_ALIGN.CENTER
            r3 = tf3.paragraphs[0].add_run(); r3.text = verdict
            r3.font.size = Pt(10); r3.font.bold = True
            r3.font.color.rgb = C_WHITE; r3.font.name = "Calibri"
            txt(sl, expl, cx + Inches(0.2), cy + Inches(0.45),
                cw - Inches(0.4), Inches(0.65), size=11, color=C_LIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE — Dockerfile + Binaires + Archives + SQL
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG_DARK)
box(sl, 0, 0, W, Inches(1.1), RGBColor(0x60, 0x20, 0xA0))
txt(sl, "05  Détail des vérifications par type (suite)",
    Inches(0.5), Inches(0.15), W, Inches(0.85), size=28, bold=True, color=C_WHITE)
divider(sl, Inches(1.15), RGBColor(0xE0, 0x70, 0xFF))

extra_types = [
    ("Dockerfile", C_ORANGE,
     [("hadolint", "Best practices CIS Docker : :latest, ADD http://, apt sans --no-install-recommends", "FAIL"),
      ("RUN curl|bash", "Téléchargement + exécution sans contrôle d'intégrité", "FAIL")]),
    ("Binaires .so .exe .elf", C_RED,
     [("FAIL automatique", "Aucun binaire précompilé dans un repo IA — implant/RAT potentiel", "FAIL"),
      ("strings + IOC", "Extraction des chaînes : URLs, /bin/sh, /etc/passwd, exec, shell", "FAIL")]),
    ("Archives .zip .tar.gz", RGBColor(0xFF, 0x80, 0x00),
     [("FAIL automatique", "Contenu non scannable sans extraction préalable", "FAIL"),
      ("Zip-slip", "Extraction vers des chemins arbitraires (../../etc/cron.d/)", "FAIL")]),
    ("SQL .sql", C_ACCENT2,
     [("xp_cmdshell", "Exécution de commandes shell depuis SQL Server", "FAIL"),
      ("DROP TABLE/DATABASE", "Instructions destructrices dans du code IA généré", "FAIL"),
      ("; --", "Pattern classique d'injection SQL", "FAIL")]),
]

for i, (name, color, checks) in enumerate(extra_types):
    col = i % 2
    row = i // 2
    cx = Inches(0.3 + col * 6.6)
    cy = Inches(1.3 + row * 3.0)

    rounded_box(sl, cx, cy, Inches(6.3), Inches(2.75), C_BG_CARD)
    s = sl.shapes.add_shape(1, cx, cy, Inches(6.3), Inches(0.5))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    txt(sl, name, cx + Inches(0.1), cy + Inches(0.05),
        Inches(6.1), Inches(0.42), size=15, bold=True, color=C_WHITE)

    for ci, (check, expl, verdict) in enumerate(checks):
        icy = cy + Inches(0.6 + ci * 0.75)
        txt(sl, f"• {check}", cx + Inches(0.15), icy,
            Inches(3.5), Inches(0.35), size=12, bold=True, color=C_RED)
        txt(sl, expl, cx + Inches(0.15), icy + Inches(0.33),
            Inches(5.8), Inches(0.35), size=11, color=C_LIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE — OUTILS (tableau synthèse)
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG_DARK)
box(sl, 0, 0, W, Inches(1.1), C_GREEN)
txt(sl, "06  Outils utilisés — Synthèse", Inches(0.5), Inches(0.15),
    W, Inches(0.85), size=32, bold=True, color=C_BG_DARK)
divider(sl, Inches(1.15), C_GREEN)

txt(sl, "⚡  Tous optionnels en dégradé — WARN si absent, FAIL uniquement sur finding actif",
    Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.45),
    size=14, italic=True, color=C_ORANGE, align=PP_ALIGN.CENTER)

tools_list = [
    ("Gitleaks",        "secrets/tokens",         "GitHub binary",         "Global",    C_ACCENT),
    ("detect-secrets",  "entropie élevée",         "pip install",           "Global",    C_ACCENT),
    ("ClamAV",          "malware signatures",      "apt install clamav",    "Global",    C_ORANGE),
    ("YARA",            "IOC personnalisés",        "apt install yara",      "Global",    RGBColor(0xE0,0x70,0xFF)),
    ("Bandit",          "SAST Python",             "pip install bandit",    ".py",       C_YELLOW),
    ("Semgrep",         "SAST multi-lang",         "pip install semgrep",   ".js .ts",   C_YELLOW),
    ("cppcheck",        "analyse statique C/C++",  "apt install cppcheck",  ".c .cpp",   C_ACCENT2),
    ("ShellCheck",      "lint shell",              "apt install shellcheck",".sh .bash", C_ACCENT2),
    ("hadolint",        "lint Dockerfile",         "GitHub binary",         "Dockerfile",C_GREEN),
    ("checkov",         "sécurité IaC",            "pip install checkov",   ".tf .hcl",  C_GREEN),
    ("jq",              "parsing JSON API",         "apt install jq",        "Utilitaire",C_GRAY),
]

hdrs = ["Outil", "Rôle", "Installation", "Scope"]
hx   = [Inches(0.3), Inches(2.85), Inches(6.8), Inches(10.5)]
hw   = [Inches(2.45), Inches(3.85), Inches(3.6),  Inches(2.5)]
for hi, (h, x, w) in enumerate(zip(hdrs, hx, hw)):
    s = sl.shapes.add_shape(1, x, Inches(1.72), w, Inches(0.42))
    s.fill.solid(); s.fill.fore_color.rgb = RGBColor(0x10, 0x60, 0x30)
    s.line.fill.background()
    txt(sl, h, x + Inches(0.05), Inches(1.75), w - Inches(0.1), Inches(0.36),
        size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

for ri, (name, role, install, scope, color) in enumerate(tools_list):
    ry = Inches(2.18 + ri * 0.48)
    bg = RGBColor(0x0E, 0x1E, 0x2E) if ri % 2 == 0 else C_BG_CARD
    vals = [name, role, install, scope]
    for vi, (val, x, w) in enumerate(zip(vals, hx, hw)):
        s = sl.shapes.add_shape(1, x, ry, w, Inches(0.44))
        s.fill.solid(); s.fill.fore_color.rgb = bg; s.line.fill.background()
        vcol = color if vi == 0 else C_LIGHT
        txt(sl, val, x + Inches(0.07), ry + Inches(0.05),
            w - Inches(0.1), Inches(0.36), size=11,
            bold=(vi == 0), color=vcol)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE — RÉSULTATS : ARCHIVE ZIP
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG_DARK)
box(sl, 0, 0, W, Inches(1.1), C_ACCENT)
txt(sl, "07  Résultats — Archive approuvée & Rapport Excel",
    Inches(0.5), Inches(0.15), W, Inches(0.85), size=28, bold=True, color=C_WHITE)
divider(sl, Inches(1.15))

# ZIP côté gauche
rounded_box(sl, Inches(0.3), Inches(1.3), Inches(5.8), Inches(5.95), C_BG_CARD)
s = sl.shapes.add_shape(1, Inches(0.3), Inches(1.3), Inches(5.8), Inches(0.55))
s.fill.solid(); s.fill.fore_color.rgb = C_GREEN; s.line.fill.background()
txt(sl, "📦  Archive ZIP  →  Good/",
    Inches(0.4), Inches(1.33), Inches(5.6), Inches(0.48),
    size=16, bold=True, color=C_BG_DARK)

zip_tree = """Good/
└── repo_20260615_143022.zip
    ├── src/
    │   └── [code source complet]
    ├── README.md
    └── scan_report_20260615.xlsx"""
mono_box(sl, zip_tree, Inches(0.45), Inches(2.0), Inches(5.5), Inches(1.85), size=11)

zip_pts = [
    "✔  Produite uniquement si verdict PASS",
    "✔  Nom horodaté — unicité garantie",
    "✔  .manifest_sha256.txt exclu",
    "✔  Rapport Excel inclus directement",
    "✔  Prête au transfert réseau interne",
    "✔  Optionnel : chiffrement GPG",
]
txt_lines(sl, zip_pts, Inches(0.45), Inches(3.95), Inches(5.6), Inches(3.1),
          size=13, color=C_LIGHT)

# Excel côté droit
rounded_box(sl, Inches(6.5), Inches(1.3), Inches(6.6), Inches(5.95), C_BG_CARD)
s2 = sl.shapes.add_shape(1, Inches(6.5), Inches(1.3), Inches(6.6), Inches(0.55))
s2.fill.solid(); s2.fill.fore_color.rgb = C_ACCENT2; s2.line.fill.background()
txt(sl, "📊  Rapport Excel (.xlsx)",
    Inches(6.6), Inches(1.33), Inches(6.4), Inches(0.48),
    size=16, bold=True, color=C_BG_DARK)

# Onglet 0
txt(sl, "Onglet 0 — Résumé",
    Inches(6.6), Inches(2.0), Inches(6.3), Inches(0.42),
    size=14, bold=True, color=C_ACCENT2)
tab0 = [
    ("Dépôt / Source",   "URL GitHub ou chemin local"),
    ("Date du scan",     "Horodatage ISO 8601"),
    ("Hash SHA-256",     "Empreinte globale de tous les fichiers"),
    ("Verdict",          "PASS (vert) / FAIL (rouge)"),
    ("Compteurs",        "PASS · WARN · FAIL"),
]
for i, (k, v) in enumerate(tab0):
    y0 = Inches(2.5 + i * 0.44)
    s3 = sl.shapes.add_shape(1, Inches(6.55), y0, Inches(2.6), Inches(0.4))
    s3.fill.solid(); s3.fill.fore_color.rgb = RGBColor(0x10, 0x30, 0x50); s3.line.fill.background()
    txt(sl, k, Inches(6.62), y0 + Inches(0.04), Inches(2.5), Inches(0.34),
        size=11, bold=True, color=C_ACCENT2)
    txt(sl, v, Inches(9.2), y0 + Inches(0.04), Inches(3.8), Inches(0.34),
        size=11, color=C_LIGHT)

# Onglet 1
txt(sl, "Onglet 1 — Fichiers",
    Inches(6.6), Inches(4.75), Inches(6.3), Inches(0.42),
    size=14, bold=True, color=C_YELLOW)
cols1 = ["#", "Fichier (chemin complet)", "Type", "Statut", "Message / Finding"]
col_colors = [C_GRAY, C_LIGHT, C_GRAY, C_YELLOW, C_LIGHT]
for i, (col, col_color) in enumerate(zip(cols1, col_colors)):
    s4 = sl.shapes.add_shape(1, Inches(6.55 + i * 1.28), Inches(5.22),
                              Inches(1.25), Inches(0.38))
    s4.fill.solid(); s4.fill.fore_color.rgb = RGBColor(0x30, 0x10, 0x60); s4.line.fill.background()
    txt(sl, col, Inches(6.58 + i * 1.28), Inches(5.25),
        Inches(1.22), Inches(0.32), size=10, bold=True, color=col_color)

for ri, (status, color_row) in enumerate([("FAIL", C_RED), ("WARN", C_ORANGE), ("PASS", C_GREEN)]):
    y5 = Inches(5.65 + ri * 0.44)
    for ci in range(5):
        s5 = sl.shapes.add_shape(1, Inches(6.55 + ci * 1.28), y5,
                                  Inches(1.25), Inches(0.4))
        bg5 = RGBColor(0x2A, 0x05, 0x05) if status=="FAIL" \
              else RGBColor(0x2A, 0x1A, 0x00) if status=="WARN" \
              else RGBColor(0x05, 0x20, 0x05)
        s5.fill.solid(); s5.fill.fore_color.rgb = bg5; s5.line.fill.background()
        val5 = ["1", "/path/file.py", ".py", status, "bandit:Severity:HIGH"][ci] \
               if status == "FAIL" else \
               ["2", "/path/config.sh", ".sh", status, "shellcheck absent"][ci] \
               if status == "WARN" else \
               ["3", "/path/main.py", ".py", status, "—"][ci]
        txt(sl, val5, Inches(6.58 + ci * 1.28), y5 + Inches(0.04),
            Inches(1.22), Inches(0.34), size=9,
            color=color_row if ci == 3 else C_LIGHT)

txt(sl, "Trié : FAIL → WARN → PASS  |  Filtre auto  |  Ligne 1 figée",
    Inches(6.55), Inches(7.1), Inches(6.5), Inches(0.35),
    size=10, italic=True, color=C_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE — RÈGLES DE SÉCURITÉ
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG_DARK)
box(sl, 0, 0, W, Inches(1.1), C_RED)
txt(sl, "08  Règles de sécurité absolues", Inches(0.5), Inches(0.15),
    W, Inches(0.85), size=32, bold=True, color=C_WHITE)
divider(sl, Inches(1.15), C_RED)

txt(sl, "⛔  Ces règles NE PEUVENT PAS être modifiées — elles constituent le modèle de menace",
    Inches(0.5), Inches(1.25), Inches(12.3), Inches(0.4),
    size=14, bold=True, color=C_RED, align=PP_ALIGN.CENTER)

rules = [
    ("🔒  Isolation réseau",
     "La machine de transit n'a JAMAIS accès direct au réseau d'entreprise"),
    ("➡️  Flux unidirectionnel",
     "Seul approved/ est lisible depuis l'interne — pas fetch/, logs/, quarantine/"),
    ("🔐  Quarantaine root-only",
     "chmod 700 — aucun utilisateur applicatif ne peut lire les fichiers refusés"),
    ("🚫  Pas de déplacement manuel",
     "Jamais de quarantine/ → approved/ sans re-scan complet du pipeline"),
    ("⚠️  Binaires = FAIL absolu",
     "Aucun .so/.exe/.elf dans un repo IA — zéro exception"),
    ("📦  Archives = re-scan obligatoire",
     ".zip/.tar.gz nécessitent extraction + re-scan dans un env isolé dédié"),
]

for i, (title, desc) in enumerate(rules):
    col = i % 2
    row = i // 2
    x = Inches(0.3 + col * 6.65)
    y = Inches(1.8 + row * 1.8)
    rounded_box(sl, x, y, Inches(6.35), Inches(1.65), C_BG_CARD)
    s = sl.shapes.add_shape(1, x, y, Inches(6.35), Inches(0.5))
    s.fill.solid()
    s.fill.fore_color.rgb = RGBColor(0x6A, 0x00, 0x00)
    s.line.fill.background()
    txt(sl, title, x + Inches(0.1), y + Inches(0.07),
        Inches(6.1), Inches(0.42), size=14, bold=True, color=C_RED)
    txt(sl, desc, x + Inches(0.15), y + Inches(0.62),
        Inches(6.0), Inches(0.85), size=13, color=C_LIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE — EXPLOITATION & ASSURANCE QUALITE
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG_DARK)
box(sl, 0, 0, W, Inches(1.0), RGBColor(0x00, 0x70, 0xA8))
txt(sl, "09  Exploitation & assurance qualité", Inches(0.5), Inches(0.12), W, Inches(0.8),
    size=30, bold=True, color=C_WHITE)
divider(sl, Inches(1.05))

# ── Colonne gauche : options d'exécution ──
txt(sl, "Options d'exécution", Inches(0.5), Inches(1.2), Inches(6), Inches(0.5),
    size=19, bold=True, color=C_ACCENT)
rounded_box(sl, Inches(0.5), Inches(1.75), Inches(5.9), Inches(3.1), C_BG_CARD)
opts = [
    ("--quiet / --verbose",   "Verbosité des journaux"),
    ("--min-severity",        "Seuil de blocage : low → critical"),
    ("--since COMMIT",        "Uniquement les fichiers modifiés"),
    ("--report-only",         "Ne bloque jamais, ne met rien en quarantaine"),
    ("--no-zip / --no-excel", "Rapports seuls, sans archive"),
]
for i, (flag, desc) in enumerate(opts):
    y = Inches(1.95 + i * 0.58)
    txt(sl, flag, Inches(0.72), y, Inches(2.6), Inches(0.3),
        size=12, bold=True, color=C_ACCENT2)
    txt(sl, desc, Inches(0.72), y + Inches(0.25), Inches(5.4), Inches(0.3),
        size=11, color=C_GRAY)

# ── Colonne droite : le pipeline se vérifie lui-même ──
txt(sl, "Le pipeline se vérifie lui-même", Inches(6.9), Inches(1.2), Inches(6), Inches(0.5),
    size=19, bold=True, color=C_GREEN)
rounded_box(sl, Inches(6.9), Inches(1.75), Inches(5.9), Inches(3.1), C_BG_CARD)
qa = [
    ("Suite de tests",         "51 assertions, aucun outil requis"),
    ("Corpus de règles",       "chaque finding doit viser le bon fichier"),
    ("Garde-fous",             "le code sûr ne doit jamais être signalé"),
    ("Intégration continue",   "lint · test · pins · docker"),
    ("Validation par mutation", "le défaut est réintroduit, le test doit échouer"),
]
for i, (title, desc) in enumerate(qa):
    y = Inches(1.95 + i * 0.58)
    txt(sl, title, Inches(7.12), y, Inches(3.0), Inches(0.3),
        size=12, bold=True, color=C_WHITE)
    txt(sl, desc, Inches(7.12), y + Inches(0.25), Inches(5.4), Inches(0.3),
        size=11, color=C_GRAY)

# ── Bas : contrôles par dépôt ──
txt(sl, "Contrôles par dépôt", Inches(0.5), Inches(5.0), Inches(6), Inches(0.4),
    size=17, bold=True, color=C_ACCENT)
txt_lines(sl, [
    ".transitignore  —  motifs gitignore, fichiers exclus de toutes les couches",
    ".transit-allow.json  —  rétrograde un FAIL connu en WARN, avec justification",
], Inches(0.5), Inches(5.45), Inches(12.3), Inches(0.9), size=12, color=C_GRAY)

txt(sl, "Une suite qu'on n'a jamais vue échouer n'apporte aucune preuve.",
    Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5),
    size=15, bold=True, italic=True, color=C_ACCENT2, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE — CONCLUSION
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
fill_bg(sl, C_BG_DARK)
box(sl, 0, 0, W, Inches(1.1), C_ACCENT2)
txt(sl, "En résumé", Inches(0.5), Inches(0.15), W, Inches(0.85),
    size=36, bold=True, color=C_BG_DARK)
divider(sl, Inches(1.15), C_ACCENT2)

recap = [
    ("🌐", "Source", "GitHub (whitelist) ou chemin local", C_ACCENT),
    ("⬇",  "Fetch",  "Clone minimal · .git supprimé · SHA-256", C_ACCENT2),
    ("🔍", "Scan",   "Couche globale + couche par type × 11 extensions", C_ORANGE),
    ("📦", "PASS",   "ZIP dans Good/ avec rapport Excel inclus", C_GREEN),
    ("🚨", "FAIL",   "Quarantaine chmod 700 + rapport JSON/HTML", C_RED),
]

for i, (icon, label, desc, color) in enumerate(recap):
    x = Inches(0.3 + i * 2.57)
    rounded_box(sl, x, Inches(1.45), Inches(2.45), Inches(4.5), C_BG_CARD)
    s = sl.shapes.add_shape(5, x + Inches(0.78), Inches(1.6),
                            Inches(0.9), Inches(0.8))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.adjustments[0] = 0.5
    tf5 = s.text_frame; tf5.paragraphs[0].alignment = PP_ALIGN.CENTER
    r5 = tf5.paragraphs[0].add_run(); r5.text = icon
    r5.font.size = Pt(22); r5.font.name = "Segoe UI Emoji"
    r5.font.color.rgb = C_WHITE

    txt(sl, label, x + Inches(0.1), Inches(2.55), Inches(2.25), Inches(0.48),
        size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
    txt(sl, desc, x + Inches(0.1), Inches(3.1), Inches(2.25), Inches(1.5),
        size=12, color=C_LIGHT, align=PP_ALIGN.CENTER)

    if i < 4:
        txt(sl, "→", Inches(2.78 + i * 2.57), Inches(3.45), Inches(0.3), Inches(0.45),
            size=22, color=C_GRAY, align=PP_ALIGN.CENTER)

txt(sl, "Pipeline entièrement en bash · Mode dégradé · Outils optionnels · Rapports horodatés",
    Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.42),
    size=14, italic=True, color=C_GRAY, align=PP_ALIGN.CENTER)

txt(sl, "github.com/Gaillotte/Claude  —  branche claude/vigilant-carson-f8twy0",
    Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.42),
    size=11, color=RGBColor(0x44, 0x55, 0x66), align=PP_ALIGN.CENTER)

# ── Sauvegarde ────────────────────────────────────────────────────────────────
pptx_path = "/home/user/Claude/AI_Transit_Pipeline_Slides.pptx"
prs.save(pptx_path)
print(f"PPTX généré : {pptx_path}")
